import os
import csv
from loguru import logger
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from litellm import completion
import pandas as pd
import json
from tenacity import retry, stop_after_attempt, wait_fixed, retry_if_exception_type
import time
from datetime import datetime

# --- Configuration ---
# Define the expected schema for extracted Q&A
QA_SCHEMA = ["question", "thought", "answer", "model"]
METADATA_FIELDS = ["filename", "page_number", "slide_index", "timestamp"]

# Prompt template for LLM extraction
# This is a basic template. It can be made more sophisticated.
PROMPT_TEMPLATE = """
你是一個專業的問答提取助手。你的任務是從提供的文本中，精確地提取出{num_qa_str}個問答對。
**所有提取的資訊，包括問題、思考過程和答案，都必須嚴格地基於提供的「文本」內容，不得引入任何外部知識或臆測。**
每個問答對必須包含以下三個鍵：'question', 'thought', 'answer'。

**提取規則：**
1.  **Question (問題):** 根據提供的「文本」內容，提出一個清晰、簡潔的問題。這個問題應該是使用者可能會提出的，並且其答案可以直接從文本中找到或推斷。**請注意，問題中不得提及檔案名稱或任何與來源文件相關的資訊。**
2.  **Thought (思考過程):** 詳細描述你是如何從「文本」中推導出答案的。這包括：
    *   你識別的關鍵詞或短語。
    *   文本中支持答案的具體句子或段落。
    *   你的邏輯推理過程，例如如何處理歧義、如何綜合多個信息點來形成答案。
    *   如果文本中沒有直接答案，請說明你是如何判斷的。
    *   **當提及文本中的特定文件或來源時，請務必輸出不含副檔名的完整檔案名稱，例如：`在{current_filename_without_ext}的「電氣安全注意事項」部分...`，以確保資訊的可追溯性。**
3.  **Answer (答案):** 提供問題的直接答案。答案必須完全基於提供的「文本」內容，**嚴禁引入外部知識**。如果文本中沒有足夠的信息來回答問題，請明確表示「文本中沒有足夠的資訊來回答此問題。」
    *   **當提及文本中的特定文件或來源時，請務必輸出不含副檔名的完整檔案名稱，例如：`在{current_filename_without_ext}的「電氣安全注意事項」部分...`，以確保資訊的可追溯性。**

**輸出格式：**
請以一個 JSON 陣列的形式返回結果，每個元素是一個問答對物件。
請確保 JSON 格式嚴格正確，並且所有鍵值對都符合上述定義。

**範例輸出格式:**
[
    {{
        "question": "機台的電氣連接應由誰執行？",
        "thought": "我在{current_filename_without_ext}的「電氣安全注意事項」部分找到了關鍵詞「電氣連接」和「執行」。具體句子是「電氣連接應由專業電工執行。」這句話直接回答了問題，因此我將其作為答案，並將找到的關鍵信息作為思考過程。",
        "answer": "在{current_filename_without_ext}中提到，電氣連接應由專業電工執行。"
    }},
    {{
        "question": "操作介質液時需要佩戴哪些防護設備？",
        "thought": "我在{current_filename_without_ext}中搜索「介質液」和「防護設備」。在「介質液安全注意事項」部分，我發現了「操作介質液時，請佩戴防護手套、安全眼鏡等個人防護設備。」這句話。這句話明確列出了所需的防護設備，因此我將其作為答案，並將搜索過程和找到的具體內容作為思考過程。",
        "answer": "在{current_filename_without_ext}中提到，操作介質液時，請佩戴防護手套、安全眼鏡等個人防護設備。"
    }},
    {{
        "question": "這份文件中是否提到了機台的保固期限？",
        "thought": "我仔細閱讀了{current_filename_without_ext}，特別是關於「保固」、「期限」、「售後服務」等相關詞彙。在文本中沒有找到任何關於機台具體保固期限的明確說明或章節。因此，判斷文本中沒有足夠的資訊來回答此問題。",
        "answer": "文本中沒有足夠的資訊來回答此問題。"
    }}
]

**文本：**
---
{extracted_text}
---
"""

# --- File Parsers ---
def parse_csv(file_path: str) -> str:
    """Parses a CSV file and returns its content as a single string."""
    content = []
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                content.append(','.join(row))
        logger.debug(f"Parsed CSV: {file_path}")
        return "\n".join(content)
    except Exception as e:
        logger.error(f"Error parsing CSV file {file_path}: {e}")
        raise

def parse_pdf(file_path: str) -> str:
    """Parses a PDF file and returns its content as a single string."""
    content = []
    try:
        with open(file_path, 'rb') as f:
            reader = PdfReader(f)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    content.append(f"--- Page {i+1} ---\n{text}")
        logger.debug(f"Parsed PDF: {file_path}")
        return "\n".join(content)
    except Exception as e:
        logger.error(f"Error parsing PDF file {file_path}: {e}")
        raise

def parse_docx(file_path: str) -> str:
    """Parses a DOCX file and returns its content as a single string."""
    content = []
    try:
        document = Document(file_path)
        for para in document.paragraphs:
            if para.text:
                content.append(para.text)
        logger.debug(f"Parsed DOCX: {file_path}")
        return "\n".join(content)
    except Exception as e:
        logger.error(f"Error parsing DOCX file {file_path}: {e}")
        raise

def parse_pptx(file_path: str) -> str:
    """Parses a PPTX file and returns its content as a single string."""
    content = []
    try:
        presentation = Presentation(file_path)
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text:
                        slide_text.append(shape.text)
            if slide_text:
                slide_content = '\n'.join(slide_text)
                content.append(f"--- Slide {i+1} ---\n{slide_content}")
        logger.debug(f"Parsed PPTX: {file_path}")
        return "\n".join(content)
    except Exception as e:
        logger.error(f"Error parsing PPTX file {file_path}: {e}")
        raise

# Map file extensions to parser functions
PARSERS = {
    'csv': parse_csv,
    'pdf': parse_pdf,
    'doc': parse_docx,  # Assuming .doc will be handled by docx parser or converted
    'docx': parse_docx,
    'ppt': parse_pptx,  # Assuming .ppt will be handled by pptx parser or converted
    'pptx': parse_pptx,
}

def get_file_parser(file_extension: str):
    """Returns the appropriate parser function for a given file extension."""
    return PARSERS.get(file_extension.lower())

# --- Litellm Integration ---
@retry(stop=stop_after_attempt(3), wait=wait_fixed(2), retry=retry_if_exception_type(Exception))
def call_litellm_api(model_name: str, text_content: str, current_filename_without_ext: str, num_qa_pairs: int = None) -> list[dict]:
    """
    Calls the litellm API to extract Q&A content from the given text.
    Includes retry mechanism.
    """
    
    num_qa_str = f"{num_qa_pairs}個" if num_qa_pairs is not None else ""
    messages = [
        {"role": "user", "content": PROMPT_TEMPLATE.format(extracted_text=text_content, num_qa_str=num_qa_str, current_filename_without_ext=current_filename_without_ext)}
    ]
    logger.debug(f"Calling litellm with model: {model_name}")
    try:
        response = completion(model=model_name, messages=messages, response_format={"type": "json_object"})
        # litellm's response structure might vary, typically content is in choices[0].message.content
        response_content = response.choices[0].message.content
        logger.debug(f"Litellm raw response: {response_content}")
        qa_pairs = json.loads(response_content)

        # Add model name to each QA pair and validate schema
        validated_qa_pairs = []
        for qa in qa_pairs:
            if all(key in qa for key in ["question", "thought", "answer"]):
                qa["model"] = model_name
                validated_qa_pairs.append(qa)
            else:
                logger.warning(f"Skipping malformed QA pair from LLM: {qa}")
        return validated_qa_pairs
    except json.JSONDecodeError as e:
        logger.error(f"Litellm response was not valid JSON: {response_content[:500]}... Error: {e}")
        raise ValueError("Invalid JSON response from LLM") from e
    except Exception as e:
        logger.error(f"Error calling litellm API with model {model_name}: {e}")
        raise

# --- Data Processing and Export ---
def process_documents(input_path: str, input_format: str, litellm_model_name: str, output_path: str, export_format: str, num_qa_pairs: int = None):
    """
    Main function to process documents, extract Q&A, and export results.
    """
    all_extracted_data = []
    files_to_process = []

    if os.path.isfile(input_path):
        files_to_process.append(input_path)
    elif os.path.isdir(input_path):
        for root, _, files in os.walk(input_path):
            for file in files:
                file_ext = file.split('.')[-1].lower()
                if input_format == 'auto' or file_ext == input_format:
                    if get_file_parser(file_ext):
                        files_to_process.append(os.path.join(root, file))
                    else:
                        logger.warning(f"Skipping unsupported file type: {file} in auto mode.")
                elif file_ext == input_format: # Specific format requested
                     if get_file_parser(file_ext):
                        files_to_process.append(os.path.join(root, file))
                     else:
                        logger.warning(f"Skipping file {file} as its extension '{file_ext}' does not match requested format '{input_format}'.")
    else:
        logger.error(f"Input path is neither a file nor a directory: {input_path}")
        return

    if not files_to_process:
        logger.warning("No supported files found to process.")
        return

    for file_path in files_to_process:
        logger.info(f"Processing file: {file_path}")
        file_name = os.path.basename(file_path)
        file_ext = file_name.split('.')[-1].lower()

        parser = get_file_parser(file_ext)
        if not parser:
            logger.warning(f"No parser available for file type: {file_ext}. Skipping {file_name}")
            continue

        try:
            extracted_text = parser(file_path)
            if not extracted_text.strip():
                logger.warning(f"No text extracted from {file_name}. Skipping LLM call.")
                continue

            file_name_without_ext = os.path.splitext(file_name)[0]
            qa_pairs = call_litellm_api(litellm_model_name, extracted_text, file_name_without_ext, num_qa_pairs)

            if num_qa_pairs is not None and len(qa_pairs) > num_qa_pairs:
                logger.warning(f"Truncating {len(qa_pairs)} Q&A pairs to {num_qa_pairs} as requested.")
                qa_pairs = qa_pairs[:num_qa_pairs]

            for qa in qa_pairs:
                # Add metadata
                qa["filename"] = file_name
                qa["timestamp"] = time.time() # Unix timestamp
                # Placeholder for page_number/slide_index - requires parser modification
                # For now, these will be None unless explicitly handled by parser
                qa["page_number"] = None
                qa["slide_index"] = None
                all_extracted_data.append(qa)

        except Exception as e:
            logger.error(f"Failed to process {file_name}: {e}")
            continue

    if not all_extracted_data:
        logger.warning("No Q&A data was extracted from any documents.")
        return

    # Ensure output directory exists if output_path is a directory
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Export data
    df = pd.DataFrame(all_extracted_data)
    
    # Generate timestamp
    timestamp_str = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # Split path into directory and filename
    output_dir_from_path = os.path.dirname(output_path)
    output_filename_from_path = os.path.basename(output_path)
    
    # Split path into directory and filename
    output_dir_from_path = os.path.dirname(output_path)
    output_filename_from_path = os.path.basename(output_path)
    
    # Split filename into base and extension
    base_filename, ext = os.path.splitext(output_filename_from_path)
    
    # Construct new filename with timestamp
    if base_filename: # If there's a base filename
        # If the original extension is the same as the export format, remove it to avoid duplication
        if ext and ext[1:].lower() == export_format.lower():
             new_filename = f"{base_filename}_{timestamp_str}"
        else:
             new_filename = f"{base_filename}_{timestamp_str}"
    else: # If output_path was just a directory
        new_filename = f"output_{timestamp_str}"
        
    # Reconstruct full path with new filename and correct extension
    output_file_path = os.path.join(output_dir_from_path, f"{new_filename}.{export_format}")

    if export_format == 'jsonl':
        df.to_json(output_file_path, orient='records', lines=True, force_ascii=False)
    elif export_format == 'csv':
        df.to_csv(output_file_path, index=False, encoding='utf-8')
    elif export_format == 'parquet':
        df.to_parquet(output_file_path, index=False)
    else:
        logger.error(f"Unsupported export format: {export_format}")
        return

    logger.info(f"Successfully exported {len(all_extracted_data)} Q&A entries to {output_file_path}")

from typing import Union
from datasets import Dataset, DatasetDict

def push_to_hub(dataset_obj: Union[Dataset, DatasetDict], repo_id: str):
    """
    Pushes a Hugging Face Dataset or DatasetDict to Hugging Face Hub.
    """
    logger.info(f"Attempting to push dataset to Hugging Face Hub repository: {repo_id}")

    try:
        from huggingface_hub import get_token
        if get_token() is None:
            logger.error("You are not logged in to Hugging Face. Please run `huggingface-cli login` in your terminal.")
            return

        dataset_obj.push_to_hub(repo_id)
        logger.info(f"Successfully pushed dataset to {repo_id} on Hugging Face Hub.")

    except ImportError:
        logger.error("The 'datasets' or 'huggingface_hub' library is not installed. Please install them using 'pip install datasets hugginggingface_hub'.")
    except Exception as e:
        logger.error(f"An error occurred while pushing to Hugging Face Hub: {e}")
        raise